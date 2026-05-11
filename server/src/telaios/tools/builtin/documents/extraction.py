"""
tools/builtin/documents/extraction.py
--------------------------------------
Consolidated document text extraction.

Vendor-agnostic document extraction module.  Supports: PDF, DOCX, XLSX, PPTX,
HTML, Markdown, EML, MSG, plain text.

Usage::

    from tools.builtin.documents.extraction import extract_text

    text = await extract_text(buffer, "application/pdf")
"""

from __future__ import annotations

import io
import logging
from typing import cast

logger = logging.getLogger(__name__)


async def extract_text(
    buffer: bytes,
    mime_type: str,
    file_type: str | None = None,
) -> str:
    """
    Extract plain text from a file buffer.

    Uses ``mime_type`` first; falls back to ``file_type`` (extension) for
    generic mime types. Returns an empty string if extraction is not supported.
    """
    mime = mime_type.lower()
    ext = (file_type or "").lower()

    # PDF
    if mime == "application/pdf" or ext == "pdf":
        return _extract_pdf(buffer)

    # DOCX (Word)
    if mime in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ) or ext in ("docx", "doc"):
        return _extract_docx(buffer)

    # XLSX (Excel)
    if mime in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ) or ext in ("xlsx", "xls"):
        return _extract_xlsx(buffer)

    # PPTX (PowerPoint)
    if mime in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ) or ext in ("pptx", "ppt"):
        return _extract_pptx(buffer)

    # HTML
    if mime in ("text/html", "application/xhtml+xml") or ext in ("html", "htm", "xhtml"):
        return _extract_html(buffer)

    # Markdown
    if ext in ("md", "markdown"):
        return _extract_markdown(buffer)

    # Email messages
    if mime == "message/rfc822" or ext == "eml":
        return _extract_eml(buffer)
    if mime == "application/vnd.ms-outlook" or ext == "msg":
        return _extract_msg(buffer)

    # Plain-text formats
    if (
        mime.startswith("text/")
        or mime in ("application/json", "application/x-ndjson")
        or ext in ("txt", "csv", "json")
    ):
        return buffer.decode("utf-8", errors="replace")

    # Generic octet-stream with known text extensions
    if mime == "application/octet-stream" and ext:
        text_exts = {
            "md",
            "txt",
            "csv",
            "json",
            "ts",
            "js",
            "py",
            "yaml",
            "yml",
            "toml",
            "xml",
            "html",
            "css",
            "sh",
            "java",
            "c",
            "cpp",
            "h",
        }
        if ext in text_exts:
            return buffer.decode("utf-8", errors="replace")

    return ""


# ── Format-specific extractors ────────────────────────────────────────────────


def _extract_pdf(buffer: bytes) -> str:
    """Extract text from PDF with layout preservation and OCR fallback."""
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=buffer, filetype="pdf")
        pages: list[str] = []
        use_ocr = False

        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text")

            if len(text.strip()) < 50:
                ocr_text = _extract_pdf_page_ocr(page)
                if ocr_text:
                    text = ocr_text
                    use_ocr = True

            pages.append(f"<!-- Page {page_num} -->\n{text}")

        doc.close()

        if use_ocr:
            logger.info("PDF: used OCR for some pages (scanned document detected)")

        return "\n\n".join(pages)
    except Exception as exc:
        logger.warning("PDF extraction failed: %s", exc)
        return ""


def _extract_pdf_page_ocr(page: object) -> str:
    """OCR a single PDF page using tesseract."""
    try:
        import io as _io

        import pytesseract
        from PIL import Image

        pix = page.get_pixmap(dpi=200)  # type: ignore[attr-defined]
        img_data = pix.tobytes("png")
        img = Image.open(_io.BytesIO(img_data))

        text = pytesseract.image_to_string(img)
        return cast(str, text).strip()
    except ImportError:
        logger.debug("OCR not available (install pytesseract + pillow)")
        return ""
    except Exception as exc:
        logger.debug("OCR failed for page: %s", exc)
        return ""


def _extract_docx(buffer: bytes) -> str:
    """Extract text from DOCX (Word documents)."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(buffer))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        return "\n".join(parts)
    except Exception as exc:
        logger.warning("DOCX extraction failed: %s", exc)
        return ""


def _extract_xlsx(buffer: bytes) -> str:
    """Extract text from XLSX (Excel spreadsheets)."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(buffer), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                rows.append(",".join(cells))
            if rows:
                parts.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)
    except Exception as exc:
        logger.warning("XLSX extraction failed: %s", exc)
        return ""


def _extract_pptx(buffer: bytes) -> str:
    """Extract text from PPTX (PowerPoint presentations)."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(buffer))
        parts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text)

            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                slide_parts.append(f"[Notes]: {slide.notes_slide.notes_text_frame.text}")

            if slide_parts:
                parts.append(f"## Slide {slide_num}\n" + "\n".join(slide_parts))

        return "\n\n".join(parts)
    except ImportError:
        logger.warning("python-pptx not installed, skipping PPTX extraction")
        return ""
    except Exception as exc:
        logger.warning("PPTX extraction failed: %s", exc)
        return ""


def _extract_html(buffer: bytes) -> str:
    """Extract text from HTML, stripping scripts, styles, and navigation."""
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(buffer, "html.parser")

        for script in soup(["script", "style", "nav", "header", "footer", "aside"]):
            script.decompose()

        text = soup.get_text(separator="\n", strip=True)

        lines = [line.strip() for line in text.split("\n") if line.strip()]
        return "\n".join(lines)
    except ImportError:
        logger.warning("beautifulsoup4 not installed, skipping HTML extraction")
        return ""
    except Exception as exc:
        logger.warning("HTML extraction failed: %s", exc)
        return ""


def _extract_markdown(buffer: bytes) -> str:
    """Extract text from Markdown, optionally parsing frontmatter."""
    try:
        import frontmatter

        content = buffer.decode("utf-8", errors="replace")
        post = frontmatter.loads(content)
        return str(post.content).strip()
    except ImportError:
        content = buffer.decode("utf-8", errors="replace")
        return content.strip()
    except Exception as exc:
        logger.warning("Markdown extraction failed: %s", exc)
        return buffer.decode("utf-8", errors="replace")


def _extract_eml(buffer: bytes) -> str:
    """Extract text from EML (email) files."""
    try:
        from email import policy as email_policy
        from email.parser import BytesParser

        msg = BytesParser(policy=email_policy.default).parsebytes(buffer)

        parts: list[str] = []

        if subject := msg.get("Subject", ""):
            parts.append(f"Subject: {subject}")
        if from_addr := msg.get("From", ""):
            parts.append(f"From: {from_addr}")
        if to_addr := msg.get("To", ""):
            parts.append(f"To: {to_addr}")

        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                if content_type == "text/plain":
                    body = part.get_content()
                    break
        else:
            body = msg.get_content()

        if body:
            parts.append(f"\n{body}")

        return "\n".join(parts)
    except Exception as exc:
        logger.warning("EML extraction failed: %s", exc)
        return ""


def _extract_msg(buffer: bytes) -> str:
    """Extract text from MSG (Outlook) files."""
    try:
        import extract_msg

        msg = extract_msg.Message(io.BytesIO(buffer))

        parts: list[str] = []

        if msg.subject:
            parts.append(f"Subject: {msg.subject}")
        if msg.sender:
            parts.append(f"From: {msg.sender}")
        if msg.to:
            parts.append(f"To: {msg.to}")

        if msg.body:
            parts.append(f"\n{msg.body}")

        return "\n".join(parts)
    except ImportError:
        logger.warning("extract-msg not installed, skipping MSG extraction")
        return ""
    except Exception as exc:
        logger.warning("MSG extraction failed: %s", exc)
        return ""
